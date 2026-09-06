from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from .models import Asset, Block, ParsedDocument, WarningItem
from .security import safe_name


def unsupported_pptx_shape_types() -> set:
    names = ("CHART", "EMBEDDED_OLE_OBJECT", "LINKED_OLE_OBJECT", "OLE_CONTROL_OBJECT")
    return {value for name in names if (value := getattr(MSO_SHAPE_TYPE, name, None)) is not None}


def _table_rows(table) -> list[list[str]]:
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def parse_docx(source: Path) -> ParsedDocument:
    document = Document(source)
    title = source.stem
    blocks: list[Block] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = next((int(char) for char in style if char.isdigit()), 1)
            blocks.append(Block("heading", text, level))
            if level == 1 and title == source.stem:
                title = text
        elif "quote" in style:
            blocks.append(Block("quote", text))
        elif "list" in style:
            blocks.append(Block("list", text, 0))
        else:
            blocks.append(Block("paragraph", text))
    for table in document.tables:
        blocks.append(Block("table", rows=_table_rows(table)))
    assets: list[Asset] = []
    for index, shape in enumerate(document.inline_shapes, start=1):
        try:
            image = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            part = document.part.related_parts[image]
            extension = part.content_type.rsplit("/", 1)[-1].replace("jpeg", "jpg")
            name = f"image-{index}.{safe_name(extension, 'bin')}"
            assets.append(Asset(name, part.blob))
            blocks.append(Block("image", name, asset_name=name))
        except (AttributeError, KeyError):
            continue
    warnings = []
    if document.inline_shapes and not assets:
        warnings.append(WarningItem("DOCX_IMAGE_EXPORT_FAILED", "有一张内嵌图片未能导出。"))
    return ParsedDocument(title, "docx", blocks, assets=assets, warnings=warnings)


def parse_pptx(source: Path) -> ParsedDocument:
    presentation = Presentation(source)
    document = ParsedDocument(source.stem, "pptx")
    for number, slide in enumerate(presentation.slides, start=1):
        document.blocks.append(Block("slide", str(number)))
        shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left, shape.shape_id))
        for shape in shapes:
            if getattr(shape, "has_table", False):
                document.blocks.append(Block("table", rows=_table_rows(shape.table)))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                except (AttributeError, ValueError):
                    document.warnings.append(
                        WarningItem(
                            "PPTX_LINKED_IMAGE_UNSUPPORTED",
                            "链接或无法读取的图片未导出。",
                            f"第 {number} 页",
                        )
                    )
                    continue
                name = (
                    f"slide-{number}-image-{len(document.assets) + 1}.{safe_name(image.ext, 'bin')}"
                )
                document.assets.append(Asset(name, image.blob))
                document.blocks.append(Block("image", name, asset_name=name))
            elif getattr(shape, "has_text_frame", False):
                is_title = bool(getattr(shape, "is_placeholder", False)) and "TITLE" in str(
                    shape.placeholder_format.type
                )
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    document.blocks.append(
                        Block(
                            "heading" if is_title else "list" if paragraph.level else "paragraph",
                            text,
                            paragraph.level + 1 if is_title else paragraph.level,
                        )
                    )
                    if is_title and number == 1 and document.title == source.stem:
                        document.title = text
            elif shape.shape_type in unsupported_pptx_shape_types():
                document.warnings.append(
                    WarningItem(
                        "PPTX_OBJECT_UNSUPPORTED",
                        "图表或嵌入对象已跳过。",
                        f"第 {number} 页",
                    )
                )
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                document.blocks.extend([Block("heading", "备注", 3), Block("paragraph", notes)])
        except AttributeError:
            pass
    return document


def parse_xlsx(source: Path) -> ParsedDocument:
    workbook = load_workbook(source, data_only=False, read_only=False)
    values = load_workbook(source, data_only=True, read_only=True)
    document = ParsedDocument(source.stem, "xlsx")
    for sheet in workbook.worksheets:
        if sheet.max_row == 1 and sheet.max_column == 1 and sheet.cell(1, 1).value is None:
            continue
        rows: list[list[str]] = []
        has_formula_without_cache = False
        cached_sheet = values[sheet.title]
        for row in sheet.iter_rows():
            values_row: list[str] = []
            for cell in row:
                value = cell.value
                if isinstance(value, str) and value.startswith("="):
                    cached = cached_sheet[cell.coordinate].value
                    has_formula_without_cache |= cached is None
                    values_row.append(value)
                else:
                    values_row.append("" if value is None else str(value))
            rows.append(values_row)
        document.sheets[sheet.title] = [Block("heading", sheet.title, 1), Block("table", rows=rows)]
        if sheet.merged_cells.ranges:
            document.warnings.append(
                WarningItem(
                    "XLSX_MERGED_CELLS_FLATTENED",
                    "合并单元格已按普通单元格导出。",
                    sheet.title,
                )
            )
        if sheet._charts:
            document.warnings.append(
                WarningItem("XLSX_CHART_UNSUPPORTED", "图表未导出。", sheet.title)
            )
        if has_formula_without_cache:
            document.warnings.append(
                WarningItem(
                    "XLSX_FORMULA_CACHE_UNAVAILABLE",
                    "公式没有可用的缓存显示值。",
                    sheet.title,
                )
            )
    values.close()
    workbook.close()
    return document


def _pdf_links(page) -> list[str]:
    links: list[str] = []
    annotations = page.get("/Annots", [])
    for annotation_reference in annotations:
        annotation = annotation_reference.get_object()
        if annotation.get("/Subtype") != "/Link":
            continue
        action = annotation.get("/A")
        if not action or action.get("/S") != "/URI":
            continue
        uri = str(action.get("/URI", "")).strip()
        if uri.startswith(("https://", "http://")) and uri not in links:
            links.append(uri)
    return links


def parse_pdf(source: Path) -> ParsedDocument:
    reader = PdfReader(source)
    title = source.stem
    if reader.metadata and reader.metadata.title:
        title = str(reader.metadata.title).strip() or source.stem
    document = ParsedDocument(title, "pdf")
    if reader.is_encrypted:
        document.warnings.append(
            WarningItem("PDF_ENCRYPTED_UNSUPPORTED", "加密 PDF 无法提取文本。")
        )
        return document
    has_text = False
    for number, page in enumerate(reader.pages, start=1):
        document.blocks.append(Block("page", str(number)))
        text = (page.extract_text() or "").strip()
        if text:
            has_text = True
            for paragraph in text.split("\n\n"):
                normalized = paragraph.strip()
                if normalized:
                    document.blocks.append(Block("paragraph", normalized))
        document.blocks.extend(Block("link", link) for link in _pdf_links(page))
    if not has_text:
        document.warnings.extend(
            (
                WarningItem("PDF_TEXT_UNAVAILABLE", "PDF 未检测到可提取文本。"),
                WarningItem(
                    "PDF_OCR_REQUIRED", "该 PDF 可能是扫描件，请先使用 OCR 生成可搜索文本。"
                ),
            )
        )
    return document


def parse_txt(source: Path) -> ParsedDocument:
    content = None
    fallback_encoding = None
    for encoding in ("utf-8-sig", "utf-16", "gb18030"):
        try:
            content = source.read_text(encoding=encoding)
            fallback_encoding = encoding if encoding != "utf-8-sig" else None
            break
        except UnicodeDecodeError:
            continue
    if content is None:
        raise ValueError("TXT 文件编码无法识别。")
    document = ParsedDocument(source.stem, "txt")
    if fallback_encoding:
        document.warnings.append(
            WarningItem("TXT_ENCODING_FALLBACK", f"TXT 使用 {fallback_encoding.upper()} 解码。")
        )
    for paragraph in content.replace("\r\n", "\n").replace("\r", "\n").split("\n\n"):
        for line in paragraph.split("\n"):
            text = line.strip()
            if not text:
                continue
            if text.startswith("# "):
                document.blocks.append(Block("heading", text[2:], 1))
            elif text.startswith(("- ", "* ", "+ ")):
                document.blocks.append(Block("list", text[2:]))
            else:
                document.blocks.append(Block("paragraph", text))
    return document


def parse_source(source: Path) -> ParsedDocument:
    parsers = {
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".xlsx": parse_xlsx,
        ".pdf": parse_pdf,
        ".txt": parse_txt,
    }
    return parsers[source.suffix.lower()](source)
