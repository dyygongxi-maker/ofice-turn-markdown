import base64
import zipfile
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from office_to_markdown import adapters
from office_to_markdown.adapters import parse_pptx, unsupported_pptx_shape_types
from office_to_markdown.batch import BatchConversionService, BatchStatus, discover_sources
from office_to_markdown.models import ConversionOptions
from office_to_markdown.security import ValidationError
from office_to_markdown.service import ConversionService
from office_to_markdown.settings import SettingsStore
from office_to_markdown.visuals import PptxVisualExporter, VisualExportError, WpsVisualExporter


def test_rejects_non_office_input(tmp_path: Path) -> None:
    source = tmp_path / "notes.txt"
    source.write_text("not an office file", encoding="utf-8")
    with pytest.raises(ValidationError, match="仅支持 DOCX"):
        ConversionService().convert(source, tmp_path)


def test_default_output_path_setting_round_trips_only_existing_directories(tmp_path: Path) -> None:
    settings_file = tmp_path / "settings.json"
    output = tmp_path / "output"
    output.mkdir()
    store = SettingsStore(settings_file)

    store.save_default_output(output)

    assert store.load_default_output() == output
    output.rmdir()
    assert store.load_default_output() is None


def test_pyinstaller_entrypoint_uses_an_absolute_package_import() -> None:
    launcher = Path("src/launcher.py").read_text(encoding="utf-8")
    assert "from office_to_markdown.app import main" in launcher
    assert "from .app import main" not in launcher


def test_pyinstaller_uses_the_working_tkinter_build_environment() -> None:
    build_script = Path("scripts/build.ps1").read_text(encoding="utf-8")
    assert ".venv-ui" in build_script
    assert "PyInstaller" in build_script
    app_name_definition = (
        '$AppName = "$([char]0x5EFE)$([char]0x5338)$([char]0x8F6C)$([char]0x6362)"'
    )
    assert app_name_definition in build_script
    assert "--name OfficeToMarkdown" in build_script
    assert "Move-Item -LiteralPath $BuiltPath -Destination $ReleasePath" in build_script
    rename_command = 'Rename-Item -LiteralPath (Join-Path $ReleasePath "OfficeToMarkdown.exe")'
    assert rename_command in build_script
    assert "--add-data" in build_script
    assert "$Base\\Lib\\tkinter;tkinter" in build_script
    assert "export_pptx_visuals.ps1;office_to_markdown" in build_script
    assert "export_wps_pptx_visuals.ps1;office_to_markdown" in build_script
    assert "--collect-all lxml" in build_script
    assert "shiboken6" not in build_script
    assert '$ErrorActionPreference = "Stop"' in build_script


def test_desktop_ui_uses_tkinter() -> None:
    app_module = Path("src/office_to_markdown/app.py").read_text(encoding="utf-8")
    assert "import tkinter as tk" in app_module
    assert "from PySide6.QtWidgets import" not in app_module


def test_desktop_ui_defines_a_structured_workbench_layout() -> None:
    app_module = Path("src/office_to_markdown/app.py").read_text(encoding="utf-8")

    assert "self.root.overrideredirect(True)" in app_module
    assert '"#EFF1F5"' in app_module
    assert '"#3660F4"' in app_module
    assert '"本地处理 · 不上传文件"' in app_module
    assert '"PPTX 视觉附件"' in app_module
    assert '"待处理文件"' in app_module
    assert 'option_add("*Font"' not in app_module


def test_desktop_ui_defines_the_redesigned_operational_hierarchy() -> None:
    app_module = Path("src/office_to_markdown/app.py").read_text(encoding="utf-8")

    assert "self.root.minsize(1040, 760)" in app_module
    assert "def _draw_empty_queue" in app_module
    assert "def _draw_title_bar" in app_module
    assert "rules_height = 208" in app_module
    assert "queue_y, queue_h = 494, max(228, height - 572)" in app_module
    assert '"尚未添加文件"' in app_module
    assert "self.include_source_link = tk.BooleanVar(value=False)" in app_module
    assert "self.export_pptx_png = tk.BooleanVar(value=False)" in app_module
    assert "def _sync_option_states" in app_module
    assert "if progress:" in app_module


def test_user_visible_output_is_localized_to_chinese(tmp_path: Path) -> None:
    source = tmp_path / "brief.docx"
    Document().save(source)

    result = ConversionService().convert(source, tmp_path)

    index = (result.output_path / "index.md").read_text(encoding="utf-8")
    report = result.report_path.read_text(encoding="utf-8")
    app_module = Path("src/office_to_markdown/app.py").read_text(encoding="utf-8")
    assert "源文件格式" in index
    assert "# 转换报告" in report
    assert 'self.root.title("廾匸转换")' in app_module


def test_rejects_extension_content_mismatch_and_macros(tmp_path: Path) -> None:
    mismatch = tmp_path / "not-a-document.docx"
    with zipfile.ZipFile(mismatch, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("xl/workbook.xml", "<workbook/>")
    with pytest.raises(ValidationError, match="不匹配"):
        ConversionService().convert(mismatch, tmp_path)

    macro = tmp_path / "macro.docx"
    with zipfile.ZipFile(macro, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("word/document.xml", "<document/>")
        archive.writestr("word/vbaProject.bin", b"macro")
    with pytest.raises(ValidationError, match="不接受包含宏"):
        ConversionService().convert(macro, tmp_path)


def test_converts_docx_to_markdown(tmp_path: Path) -> None:
    source = tmp_path / "brief.docx"
    document = Document()
    document.add_heading("Project Brief", level=1)
    document.add_paragraph("A local conversion test.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Mode"
    table.cell(1, 1).text = "Local"
    document.save(source)

    result = ConversionService().convert(source, tmp_path)

    assert (result.output_path / "index.md").is_file()
    content = (result.output_path / "markdown" / "brief.md").read_text(encoding="utf-8")
    assert result.report_path == result.output_path / "reports" / "brief转换报告.md"
    assert "# Project Brief" in content
    assert "| Name | Value |" in content


def test_stores_named_markdown_and_report_in_separate_directories(tmp_path: Path) -> None:
    source = tmp_path / "知识笔记.docx"
    document = Document()
    document.add_paragraph("A local conversion test.")
    document.save(source)

    result = ConversionService().convert(source, tmp_path)

    markdown_path = result.output_path / "markdown" / "知识笔记.md"
    report_path = result.output_path / "reports" / "知识笔记转换报告.md"
    index = (result.output_path / "index.md").read_text(encoding="utf-8")
    assert markdown_path.is_file()
    assert result.report_path == report_path
    assert "[转换内容](markdown/知识笔记.md)" in index
    assert "[转换报告](reports/知识笔记转换报告.md)" in index


def test_converts_pptx_to_markdown(tmp_path: Path) -> None:
    source = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Notes"
    slide.placeholders[1].text = "First point\nSecond point"
    presentation.save(source)

    result = ConversionService().convert(source, tmp_path)

    content = (result.output_path / "markdown" / "slides.md").read_text(encoding="utf-8")
    assert result.report_path == result.output_path / "reports" / "slides转换报告.md"
    assert "## 第 1 页" in content
    assert "Quarterly Notes" in content


def test_pptx_can_export_visual_previews_with_a_local_renderer(tmp_path: Path) -> None:
    source = tmp_path / "slides.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(source)

    class FakeVisualExporter:
        calls: list[tuple[Path, Path, bool, bool]] = []

        def export(
            self, source_path: Path, output_dir: Path, export_png: bool, export_pdf: bool
        ) -> None:
            self.calls.append((source_path, output_dir, export_png, export_pdf))
            pages = output_dir / "pages"
            pages.mkdir(parents=True)
            (pages / "slide-001.png").write_bytes(b"png")
            (output_dir / "slides.pdf").write_bytes(b"pdf")

    exporter = FakeVisualExporter()
    result = ConversionService(visual_exporter=exporter).convert(
        source,
        tmp_path,
        ConversionOptions(export_pptx_png=True, export_pptx_pdf=True),
    )

    assert len(exporter.calls) == 1
    assert exporter.calls[0][0] == source
    assert exporter.calls[0][1].name == "visuals"
    assert exporter.calls[0][2:] == (True, True)
    assert (result.output_path / "visuals" / "pages" / "slide-001.png").is_file()
    assert (result.output_path / "visuals" / "slides.pdf").is_file()
    index = (result.output_path / "index.md").read_text(encoding="utf-8")
    assert "[每页预览图](visuals/pages/)" in index
    assert "[版式 PDF](visuals/slides.pdf)" in index


def test_pptx_visual_export_starts_powershell_in_sta_mode(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "slides.pptx"
    source.write_bytes(b"pptx")
    output = tmp_path / "visuals"
    calls: list[list[str]] = []

    class CompletedProcess:
        returncode = 0

    def fake_run(command: list[str], **_kwargs: object) -> CompletedProcess:
        calls.append(command)
        pages = output / "pages"
        pages.mkdir(parents=True)
        (pages / "slide-001.png").write_bytes(b"png")
        (output / "slides.pdf").write_bytes(b"pdf")
        return CompletedProcess()

    monkeypatch.setattr("office_to_markdown.visuals.subprocess.run", fake_run)

    PptxVisualExporter().export(source, output, export_png=True, export_pdf=True)

    assert calls[0][:3] == ["powershell.exe", "-Sta", "-NoProfile"]


def test_wps_visual_export_uses_the_wps_automation_script(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "slides.pptx"
    source.write_bytes(b"pptx")
    output = tmp_path / "visuals"
    calls: list[list[str]] = []

    class CompletedProcess:
        returncode = 0

    def fake_run(command: list[str], **_kwargs: object) -> CompletedProcess:
        calls.append(command)
        pages = output / "pages"
        pages.mkdir(parents=True)
        (pages / "slide-001.png").write_bytes(b"png")
        (output / "slides.pdf").write_bytes(b"pdf")
        return CompletedProcess()

    monkeypatch.setattr("office_to_markdown.visuals.subprocess.run", fake_run)

    WpsVisualExporter().export(source, output, export_png=True, export_pdf=True)

    decoded = base64.b64decode(calls[0][-1]).decode("utf-16-le")
    assert calls[0][:3] == ["powershell.exe", "-Sta", "-NoProfile"]
    assert "export_wps_pptx_visuals.ps1" in decoded


def test_default_visual_exporter_falls_back_to_powerpoint_after_wps_failure(
    tmp_path: Path,
) -> None:
    source = tmp_path / "slides.pptx"
    source.write_bytes(b"pptx")
    output = tmp_path / "visuals"
    calls: list[str] = []

    class FailingWpsExporter:
        def export(self, *_args: object) -> None:
            calls.append("wps")
            raise VisualExportError("WPS unavailable")

    class WorkingPowerPointExporter:
        def export(self, _source: Path, output_dir: Path, *_args: object) -> None:
            calls.append("powerpoint")
            pages = output_dir / "pages"
            pages.mkdir(parents=True)
            (pages / "slide-001.png").write_bytes(b"png")
            (output_dir / "slides.pdf").write_bytes(b"pdf")

    PptxVisualExporter(
        wps_exporter=FailingWpsExporter(),
        powerpoint_exporter=WorkingPowerPointExporter(),
    ).export(source, output, export_png=True, export_pdf=True)

    assert calls == ["wps", "powerpoint"]
    assert (output / "slides.pdf").is_file()


def test_pptx_visual_export_failure_keeps_markdown_and_reports_a_warning(tmp_path: Path) -> None:
    source = tmp_path / "slides.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(source)

    class UnavailableVisualExporter:
        def export(
            self, source_path: Path, output_dir: Path, export_png: bool, export_pdf: bool
        ) -> None:
            raise VisualExportError("unavailable")

    result = ConversionService(visual_exporter=UnavailableVisualExporter()).convert(
        source,
        tmp_path,
        ConversionOptions(export_pptx_png=True, export_pptx_pdf=True),
    )

    index = (result.output_path / "index.md").read_text(encoding="utf-8")
    report = result.report_path.read_text(encoding="utf-8")
    assert (result.output_path / "markdown" / "slides.md").is_file()
    assert not (result.output_path / "visuals").exists()
    assert "PPTX_VISUAL_EXPORT_FAILED" in report
    assert "每页预览图" not in index
    assert "版式 PDF" not in index


def test_pptx_adapter_uses_available_ole_shape_enums() -> None:
    unsupported_types = unsupported_pptx_shape_types()

    assert unsupported_types


def test_pptx_linked_picture_is_skipped_with_a_warning(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    class LinkedPicture:
        shape_type = MSO_SHAPE_TYPE.PICTURE
        top = 0
        left = 0
        shape_id = 1
        has_table = False
        has_text_frame = False

        @property
        def image(self):
            raise ValueError("no embedded image")

    class Slide:
        shapes = [LinkedPicture()]

        @property
        def notes_slide(self):
            raise AttributeError

    class FakePresentation:
        slides = [Slide()]

    monkeypatch.setattr(adapters, "Presentation", lambda _: FakePresentation())

    document = parse_pptx(tmp_path / "linked-picture.pptx")

    assert not document.assets
    assert [warning.code for warning in document.warnings] == ["PPTX_LINKED_IMAGE_UNSUPPORTED"]


def test_converts_xlsx_and_warns_for_missing_formula_cache(tmp_path: Path) -> None:
    source = tmp_path / "budget.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Item", "Amount"])
    sheet.append(["Hosting", 20])
    sheet["B3"] = "=SUM(B2:B2)"
    workbook.save(source)

    result = ConversionService().convert(source, tmp_path)

    workbook_markdown = (result.output_path / "markdown" / "budget.md").read_text(encoding="utf-8")
    sheet_markdown = (result.output_path / "markdown" / "sheets" / "Budget.md").read_text(
        encoding="utf-8"
    )
    report = result.report_path.read_text(encoding="utf-8")
    assert "[Budget](sheets/Budget.md)" in workbook_markdown
    assert result.report_path == result.output_path / "reports" / "budget转换报告.md"
    assert "=SUM(B2:B2)" in sheet_markdown
    assert "XLSX_FORMULA_CACHE_UNAVAILABLE" in report


def test_does_not_overwrite_existing_output(tmp_path: Path) -> None:
    source = tmp_path / "duplicate.docx"
    Document().save(source)
    (tmp_path / "duplicate-markdown").mkdir()
    with pytest.raises(ValidationError, match="已存在"):
        ConversionService().convert(source, tmp_path)


def test_obsidian_mode_writes_safe_frontmatter_and_relative_source_link(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    source_dir = vault / "inbox"
    source_dir.mkdir()
    source = source_dir / "brief.docx"
    Document().save(source)

    result = ConversionService().convert(
        source,
        vault,
        ConversionOptions(
            obsidian_mode=True,
            include_frontmatter=True,
            include_source_link=True,
            source_link_root=vault,
            tags=("office-import", "inbox"),
        ),
    )

    index = (result.output_path / "index.md").read_text(encoding="utf-8")
    assert index.startswith('---\nsource_file: "brief.docx"')
    assert "tags:\n  - office-import\n  - inbox" in index
    assert "[原文件](../inbox/brief.docx)" in index
    assert str(vault) not in index


def test_obsidian_mode_rejects_invalid_tags_and_outside_source_root(tmp_path: Path) -> None:
    source = tmp_path / "brief.docx"
    Document().save(source)
    vault = tmp_path / "vault"
    vault.mkdir()
    with pytest.raises(ValidationError, match="标签"):
        ConversionService().convert(source, vault, ConversionOptions(tags=("bad:tag",)))
    with pytest.raises(ValidationError, match="归档根目录"):
        ConversionService().convert(
            source,
            vault,
            ConversionOptions(include_source_link=True, source_link_root=vault),
        )


def test_obsidian_source_link_is_relative_to_nested_output_directory(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    source_dir = vault / "inbox"
    output_dir = vault / "imports" / "today"
    source_dir.mkdir(parents=True)
    output_dir.mkdir(parents=True)
    source = source_dir / "brief.docx"
    Document().save(source)

    result = ConversionService().convert(
        source,
        output_dir,
        ConversionOptions(
            obsidian_mode=True,
            include_source_link=True,
            source_link_root=vault,
        ),
    )

    index = (result.output_path / "index.md").read_text(encoding="utf-8")
    assert "[原文件](../../../inbox/brief.docx)" in index


def test_batch_converts_remaining_sources_after_unexpected_parser_error(tmp_path: Path) -> None:
    source = tmp_path / "valid.docx"
    Document().save(source)

    class UnstableConverter:
        calls = 0

        def convert(self, source_path, output_parent, options):
            self.calls += 1
            if self.calls == 1:
                raise RuntimeError("parser failure")
            return ConversionService().convert(source_path, output_parent, options)

    result = BatchConversionService(UnstableConverter()).convert((source, source), tmp_path)

    assert result.items[0].status == BatchStatus.FAILED
    assert result.items[0].error_code == "CONVERSION_FAILED"
    assert result.items[1].status == BatchStatus.SUCCESS


def test_batch_emits_running_and_finished_events(tmp_path: Path) -> None:
    source = tmp_path / "valid.docx"
    Document().save(source)
    events = []

    BatchConversionService().convert((source,), tmp_path, on_item=events.append)

    assert [item.status for item in events] == [BatchStatus.RUNNING, BatchStatus.SUCCESS]


def test_batch_discovery_and_partial_failure_do_not_stop_other_files(tmp_path: Path) -> None:
    source_dir = tmp_path / "sources"
    source_dir.mkdir()
    valid = source_dir / "valid.docx"
    Document().save(valid)
    invalid = source_dir / "broken.docx"
    invalid.write_text("not OOXML", encoding="utf-8")
    (source_dir / "notes.txt").write_text("ignore", encoding="utf-8")

    assert discover_sources(source_dir, recursive=False) == (invalid, valid)
    result = BatchConversionService().convert((invalid, valid), tmp_path)

    assert result.count(BatchStatus.FAILED) == 1
    assert result.count(BatchStatus.SUCCESS) == 1
    assert (tmp_path / "valid-markdown" / "index.md").is_file()


def test_batch_skips_existing_output_and_honors_cancel_before_next_item(tmp_path: Path) -> None:
    first = tmp_path / "first.docx"
    second = tmp_path / "second.docx"
    Document().save(first)
    Document().save(second)
    (tmp_path / "first-markdown").mkdir()
    service = BatchConversionService()
    service.cancel()

    result = service.convert((first, second), tmp_path)

    assert all(item.status == BatchStatus.CANCELLED for item in result.items)
